"""
Insert a "From Research to Production: AWS Deployment Architecture" slide
into Deep_learning_presentation_v3.pptx, positioned just before the
Conclusion. Re-uses the deck's existing palette (deep blue 023E5C,
teal 1C7293, mint 02C39A) and the 13.33 x 7.5 in slide canvas.

Run: python aws/build_arch_slide.py
Output: Deep_learning_presentation_v4.pptx (alongside the v3 source).
"""

import copy
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


PROJECT = Path("/Users/jamesoon/Library/Mobile Documents/com~apple~CloudDocs/Desktop/PROJECTS/SUTD/MSTR-DAIE/DeepLearning/Project")
SRC     = PROJECT / "Deep_learning_presentation_v3.pptx"
DST     = PROJECT / "Deep_learning_presentation_v4.pptx"

NAVY = RGBColor(0x02, 0x3E, 0x5C)
TEAL = RGBColor(0x1C, 0x72, 0x93)
MINT = RGBColor(0x02, 0xC3, 0x9A)
ICE  = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK   = RGBColor(0x21, 0x29, 0x5C)
MUTED = RGBColor(0x5A, 0x6B, 0x7C)


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_box(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def add_arrow(slide, x, y, w, h, color=TEAL):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def add_pipeline_node(slide, x, y, w, h, title, sub):
    add_box(slide, x, y, w, h, fill=WHITE, line=TEAL)
    add_text(slide, x, y + 0.05, w, 0.35, title, size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + 0.40, w, h - 0.40, sub, size=9, color=MUTED, align=PP_ALIGN.CENTER)


def insert_slide_at(prs, src_idx, dst_idx):
    """Duplicate slide at src_idx into dst_idx position (returns new slide).

    python-pptx has no built-in 'insert at position' so we move via XML.
    """
    # Clone layout from any existing blank-ish slide
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[5]
    new = prs.slides.add_slide(blank_layout)
    # Reorder XML: move newly added slide id-list entry to dst_idx
    sldIdLst = prs.slides._sldIdLst
    children = list(sldIdLst)
    moved = children.pop(-1)  # the just-added slide is last
    sldIdLst.insert(dst_idx, moved)
    return new


def build_arch_slide(slide):
    # Background — soft cream
    add_box(slide, 0, 0, 13.33, 7.5, fill=RGBColor(0xF7, 0xF8, 0xFA), shape=MSO_SHAPE.RECTANGLE)

    # Header strip
    add_box(slide, 0, 0, 13.33, 1.05, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
    add_text(slide, 0.5, 0.18, 12.3, 0.45,
             "From Research to Production — AWS Deployment Architecture",
             size=26, bold=True, color=WHITE, font="Calibri")
    add_text(slide, 0.5, 0.65, 12.3, 0.35,
             "Live MedMCQA exam platform   ·   dl.mdaie-sutd.fit   ·   serverless on ap-southeast-1",
             size=12, color=ICE, font="Calibri")

    # Mint accent bar
    add_box(slide, 0, 1.05, 13.33, 0.06, fill=MINT, shape=MSO_SHAPE.RECTANGLE)

    # Section label
    add_text(slide, 0.5, 1.25, 12.3, 0.30,
             "REQUEST PATH",
             size=10, bold=True, color=TEAL, font="Calibri")

    # Pipeline row — 5 nodes + 4 arrows, fills 0.5"–12.83" horizontally (12.33" usable)
    # Node width ~ 2.05", arrow width ~ 0.45", gaps 0.05"
    NODE_W, NODE_H = 2.10, 1.35
    ARROW_W, ARROW_H = 0.55, 0.45
    Y0 = 1.65
    cur_x = 0.5

    nodes = [
        ("Browser SPA",
         "React + Vite\nCloudFront / S3\nOAC, 404→index.html"),
        ("API Gateway v2",
         "HTTP API\nCognito JWT auth\nCORS for dl.mdaie-sutd.fit"),
        ("Lambda + DynamoDB",
         "questions · submit\nsession-trigger\nthree DDB tables"),
        ("Step Functions",
         "Map (concurrency=10)\n→ grade · → aggregate\nretry: 3× backoff 2"),
        ("HF Inference",
         "Gemma-3-4B-it MedMCQA LoRA\n(jamezoon/...)\nDeBERTa-v3-large grader"),
    ]
    for i, (title, sub) in enumerate(nodes):
        add_pipeline_node(slide, cur_x, Y0, NODE_W, NODE_H, title, sub)
        cur_x += NODE_W
        if i < len(nodes) - 1:
            add_arrow(slide, cur_x, Y0 + (NODE_H - ARROW_H) / 2, ARROW_W, ARROW_H, color=TEAL)
            cur_x += ARROW_W

    # Two callout panels below pipeline
    add_text(slide, 0.5, 3.30, 12.3, 0.30,
             "DEMO CONTROLS  ·  RESEARCH-TO-PRODUCTION TIE-BACK",
             size=10, bold=True, color=TEAL, font="Calibri")

    # Left callout — Demo mode
    add_box(slide, 0.5, 3.65, 6.0, 2.55, fill=WHITE, line=TEAL)
    add_box(slide, 0.5, 3.65, 6.0, 0.50, fill=TEAL, shape=MSO_SHAPE.RECTANGLE)
    add_text(slide, 0.65, 3.72, 5.7, 0.40,
             "Demo Mode (SSM-driven, no redeploy)", size=14, bold=True, color=WHITE, font="Calibri")
    demo_lines = [
        ("•  Toggle: ", "aws ssm put-parameter --name /medmcqa/demo_mode --value 1"),
        ("•  Effect: ", "grade Lambda returns canned explanations from demo_explanations.json"),
        ("•  Trade-off: ", "sub-second, deterministic, $0 HF cost — Step Functions Map still visible"),
        ("•  Fallback: ", "DEMO_MODE=0 → live HF call to Gemma-3-4B-it MedMCQA LoRA"),
    ]
    y = 4.30
    for label, body in demo_lines:
        b = slide.shapes.add_textbox(Inches(0.65), Inches(y), Inches(5.7), Inches(0.42))
        tf = b.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.0)
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = label; r1.font.bold = True; r1.font.size = Pt(11); r1.font.color.rgb = NAVY; r1.font.name = "Calibri"
        r2 = p.add_run(); r2.text = body;  r2.font.size = Pt(11); r2.font.color.rgb = INK;  r2.font.name = "Consolas" if "aws ssm" in body else "Calibri"
        y += 0.46

    # Right callout — Same artifact
    add_box(slide, 6.83, 3.65, 6.0, 2.55, fill=WHITE, line=MINT)
    add_box(slide, 6.83, 3.65, 6.0, 0.50, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
    add_text(slide, 6.98, 3.72, 5.7, 0.40,
             "Deployed Model = Reported Model", size=14, bold=True, color=WHITE, font="Calibri")

    # Inner panel: model id + scorecard
    add_text(slide, 6.98, 4.30, 5.7, 0.30,
             "jamezoon / gemma-3-4b-it-medmcqa-lora",
             size=11, bold=True, color=NAVY, font="Consolas")

    # 3-stat row
    stats = [
        ("45.4 %", "MedMCQA-dev", "4,183 Q"),
        ("48.1 %", "MCAT (7 sets)", "1,609 Q"),
        ("LoRA r=16", "α=32, dropout 0.05", "q/k/v/o_proj"),
    ]
    sx = 6.98
    sw = 1.85
    for big, mid, small in stats:
        add_box(slide, sx, 4.65, sw, 1.40, fill=RGBColor(0xF0, 0xF6, 0xF8), line=ICE)
        add_text(slide, sx, 4.72, sw, 0.55, big, size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font="Calibri")
        add_text(slide, sx, 5.30, sw, 0.30, mid, size=10, bold=True, color=TEAL, align=PP_ALIGN.CENTER, font="Calibri")
        add_text(slide, sx, 5.62, sw, 0.30, small, size=9, color=MUTED, align=PP_ALIGN.CENTER, font="Calibri")
        sx += sw + 0.05

    # Caption tying back to the report
    add_text(slide, 0.5, 6.40, 12.3, 0.45,
             "The HF Inference endpoint behind the live exam app serves the same LoRA adapter "
             "evaluated in the report — the platform is the productionised research artifact, not a separate exhibit.",
             size=11, color=MUTED, align=PP_ALIGN.CENTER, font="Calibri")

    # Footer strip
    add_box(slide, 0, 7.05, 13.33, 0.45, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
    add_text(slide, 0.5, 7.12, 12.3, 0.30,
             "61.502 Deep Learning for Enterprise · SUTD 2026  |  AWS CDK · Cognito · Step Functions · HF Inference  |  github → /aws",
             size=9, color=ICE, align=PP_ALIGN.CENTER, font="Calibri")


def main():
    prs = Presentation(SRC)
    print(f"Source has {len(prs.slides)} slides ({Emu(prs.slide_width).inches:.2f} x {Emu(prs.slide_height).inches:.2f} in).")

    # Insert before the Conclusion slide.
    # Slides indexed 0..14; conclusion is slide 15 (index 14).
    new = insert_slide_at(prs, src_idx=0, dst_idx=14)
    build_arch_slide(new)

    prs.save(DST)
    print(f"Wrote {DST} with {len(prs.slides)} slides; new slide is at position {15} (1-indexed).")


if __name__ == "__main__":
    main()
